"""Generate warm-design Business Model Canvas PowerPoint for designtuntas.id"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

CREAM = RGBColor(0xF7, 0xF0, 0xE8)
TERRACOTTA = RGBColor(0xC4, 0x6B, 0x4A)
COPPER = RGBColor(0xA8, 0x56, 0x3A)
INK = RGBColor(0x3D, 0x2C, 0x29)
MUTED = RGBColor(0x7A, 0x65, 0x5C)
GOLD = RGBColor(0xB8, 0x95, 0x6C)
WHITE = RGBColor(0xFF, 0xFB, 0xF7)
CARD_A = RGBColor(0xF3, 0xE8, 0xDC)
CARD_B = RGBColor(0xED, 0xDF, 0xD0)
CARD_C = RGBColor(0xE8, 0xD4, 0xC0)
VP = RGBColor(0xF8, 0xE4, 0xD4)
COST = RGBColor(0xEA, 0xDC, 0xCB)
REV = RGBColor(0xF0, 0xD9, 0xC4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_run(run, size=12, bold=False, color=INK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=CREAM):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def add_rect(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, size=12, bold=False, color=INK, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_bullets(slide, x, y, w, h, title, items, fill=CARD_A, title_color=TERRACOTTA, body_size=9):
    add_rect(slide, x, y, w, h, fill)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TERRACOTTA
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(
        x + Inches(0.18), y + Inches(0.12), w - Inches(0.28), Inches(0.32)
    )
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run(run, size=11, bold=True, color=title_color)

    body = slide.shapes.add_textbox(
        x + Inches(0.18), y + Inches(0.42), w - Inches(0.28), h - Inches(0.52)
    )
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = "• " + item
        set_run(run, size=body_size, color=INK)


# ----- Slide 1: Cover -----
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1, CREAM)

panel = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.8), prs.slide_height)
panel.fill.solid()
panel.fill.fore_color.rgb = TERRACOTTA
panel.line.fill.background()

soft = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(4.8), Inches(4), Inches(4))
soft.fill.solid()
soft.fill.fore_color.rgb = COPPER
soft.line.fill.background()

add_textbox(s1, Inches(0.55), Inches(2.0), Inches(3.8), Inches(0.7), "designtuntas.id", 28, True, WHITE, "Georgia")
add_textbox(s1, Inches(0.55), Inches(2.85), Inches(3.8), Inches(0.5), "Business Model Canvas", 18, False, WHITE)
add_textbox(
    s1,
    Inches(0.55),
    Inches(3.6),
    Inches(3.8),
    Inches(1.2),
    "Landasan bisnis yang hangat,\nsederhana, dan siap tumbuh.",
    13,
    False,
    RGBColor(0xFF, 0xE8, 0xDE),
)

add_textbox(
    s1,
    Inches(5.5),
    Inches(2.0),
    Inches(7.2),
    Inches(1.4),
    "Desain & dokumen,\nsampai tuntas.",
    36,
    True,
    INK,
    "Georgia",
)
add_textbox(
    s1,
    Inches(5.5),
    Inches(3.8),
    Inches(7.2),
    Inches(1.6),
    "Resume CV  ·  Konsultasi Skripsi  ·  Design Visual  ·  Design 3D\n\n"
    "Kebijakan biaya: versi gratis dulu.\nBeralih berbayar setelah ada transaksi nyata.",
    14,
    False,
    MUTED,
)
add_textbox(
    s1,
    Inches(5.5),
    Inches(6.5),
    Inches(7.2),
    Inches(0.4),
    "Pondasi bisnis  ·  Fase 1: landing + Tuti AI  ·  WhatsApp: 088901178816",
    11,
    False,
    GOLD,
)

# ----- Slide 2: Full canvas -----
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2, CREAM)
add_textbox(s2, Inches(0.35), Inches(0.22), Inches(10), Inches(0.4), "Business Model Canvas — designtuntas.id", 22, True, INK, "Georgia")
add_textbox(s2, Inches(0.35), Inches(0.62), Inches(12), Inches(0.28), "Satu pintu untuk CV, skripsi, visual, dan 3D — sampai tuntas.", 11, False, MUTED)

gap = Inches(0.1)
left = Inches(0.3)
top = Inches(1.0)
col_w = Inches(2.38)
vp_w = Inches(2.75)
h_stack = Inches(2.28)
h_bot = Inches(1.5)

x0 = left
x1 = x0 + col_w + gap
x2 = x1 + col_w + gap
x3 = x2 + vp_w + gap
x4 = x3 + col_w + gap
y0 = top
y1 = y0 + h_stack + gap
y2 = y1 + h_stack + gap

add_bullets(
    s2, x0, y0, col_w, h_stack * 2 + gap,
    "Key Partners",
    [
        "OpenRouter (Tuti AI)",
        "Vercel Hobby (hosting)",
        "Nanti: Midtrans / Xendit",
        "Nanti: Supabase / storage",
        "Freelancer (opsional)",
        "Kampus / komunitas (referral)",
    ],
    CARD_B,
)

add_bullets(
    s2, x1, y0, col_w, h_stack,
    "Key Activities",
    [
        "Produksi CV, skripsi, visual, 3D",
        "Brief & manajemen revisi",
        "Konten IG / TikTok",
        "Jaga website + Tuti AI",
        "Support WhatsApp",
    ],
    CARD_A,
)

add_bullets(
    s2, x1, y1, col_w, h_stack,
    "Key Resources",
    [
        "Brand designtuntas.id",
        "Skill inti 4 layanan",
        "Web + Tuti + WhatsApp",
        "Portofolio & template",
        "Waktu founder / tim kecil",
    ],
    CARD_C,
)

add_bullets(
    s2, x2, y0, vp_w, h_stack * 2 + gap,
    "Value Propositions",
    [
        "Satu pintu: CV · skripsi · visual · 3D",
        "Sampai tuntas (revisi terarah)",
        "Proses: chat → brief → revisi → selesai",
        "WhatsApp + Tuti AI (FAQ cepat)",
        "Bahasa & konteks lokal Indonesia",
        "Harga transparan bertahap",
    ],
    VP,
    COPPER,
    10,
)

add_bullets(
    s2, x3, y0, col_w, h_stack,
    "Customer Relationships",
    [
        "Tuti AI: info mandiri",
        "WhatsApp: bantuan personal",
        "Update status & revisi jelas",
        "Nanti: tracking & akun",
        "Upsell lintas layanan",
    ],
    CARD_A,
)

add_bullets(
    s2, x3, y1, col_w, h_stack,
    "Channels",
    [
        "Website / Vercel",
        "WhatsApp (konversi utama)",
        "Instagram & TikTok",
        "Tuti AI di website",
        "Nanti: form order & payment",
    ],
    CARD_C,
)

add_bullets(
    s2, x4, y0, col_w, h_stack * 2 + gap,
    "Customer Segments",
    [
        "Mahasiswa (skripsi / dokumen)",
        "Fresh graduate & pencari kerja",
        "UMKM / personal brand",
        "Klien design 3D",
        "Prioritas awal: CV + skripsi",
    ],
    CARD_B,
)

left_bottom_w = (x2 + vp_w) - x0
right_bottom_w = (x4 + col_w) - x3

add_bullets(
    s2, x0, y2, left_bottom_w, h_bot,
    "Cost Structure",
    [
        "Gratis dulu: Vercel Hobby, OpenRouter :free, tools memadai",
        "Bayar setelah transaksi: fee gateway, domain .id, upgrade bila perlu",
    ],
    COST,
    body_size=10,
)

add_bullets(
    s2, x3, y2, right_bottom_w, h_bot,
    "Revenue Streams",
    [
        "Paket CV · Skripsi (bab/paket) · Design visual · Design 3D",
        "Opsional nanti: revisi ekstra / retainer — belum andalkan iklan/SaaS",
    ],
    REV,
    body_size=10,
)

# ----- Slide 3: Relationships explained + cost policy -----
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3, CREAM)
add_textbox(s3, Inches(0.5), Inches(0.35), Inches(12), Inches(0.5), "Hubungan pelanggan & kebijakan biaya", 24, True, INK, "Georgia")

add_bullets(
    s3,
    Inches(0.5),
    Inches(1.2),
    Inches(5.9),
    Inches(3.2),
    "Customer Relationships (jelas)",
    [
        "Tuti AI: jawab pertanyaan umum 24/7 (self-serve)",
        "WhatsApp: brief, revisi, order (bantuan personal)",
        "Alur: tanya Tuti → lanjut WA saat siap order",
        "Tidak ada refund — disepakati sejak awal dengan klien",
        "Revisi selama proses hingga valid; setelah valid berbayar",
    ],
    CARD_A,
    body_size=12,
)

add_bullets(
    s3,
    Inches(6.7),
    Inches(1.2),
    Inches(5.9),
    Inches(3.2),
    "Cost — tetap gratis dulu",
    [
        "Fase 1–3: pakai free tier selama memungkinkan",
        "Vercel Hobby + OpenRouter model free",
        "Payment sandbox sampai ada transaksi nyata",
        "Fee gateway mulai transaksi production pertama",
        "Upgrade DB/AI hanya jika volume sudah besar",
    ],
    CARD_B,
    body_size=12,
)

add_rect(s3, Inches(0.5), Inches(4.7), Inches(12.1), Inches(2.0), VP)
add_textbox(
    s3,
    Inches(0.8),
    Inches(4.95),
    Inches(11.5),
    Inches(1.5),
    "Keputusan terkunci\n"
    "Pertahankan versi gratis. Beralih berbayar hanya setelah ada transaksi nyata.\n"
    "Asisten AI: Tuti  ·  WA: 088901178816  ·  Tagline: Siap menyelesaikan masalah Anda sampai tuntas.",
    14,
    False,
    INK,
)

# ----- Slide 4: Roadmap -----
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4, CREAM)
add_textbox(s4, Inches(0.5), Inches(0.35), Inches(12), Inches(0.5), "Roadmap pertumbuhan", 24, True, INK, "Georgia")
add_textbox(s4, Inches(0.5), Inches(0.9), Inches(12), Inches(0.35), "Dari landing hangat ke platform penuh — tanpa terburu bayar.", 13, False, MUTED)

phases = [
    ("Fase 1", "Landing + Tuti AI", "Website, WhatsApp, OpenRouter free\nVercel Hobby — GRATIS", TERRACOTTA),
    ("Fase 2", "Sistem order", "Form order, notifikasi, tracking\nFree tier dulu", COPPER),
    ("Fase 3", "Platform penuh", "Akun, dashboard, payment\nSandbox dulu; bayar saat transaksi", GOLD),
]

for i, (phase, title, detail, accent) in enumerate(phases):
    x = Inches(0.5) + i * Inches(4.15)
    card = add_rect(s4, x, Inches(1.7), Inches(3.9), Inches(4.2), CARD_A)
    top_bar = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.7), Inches(3.9), Inches(0.7))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent
    top_bar.line.fill.background()
    add_textbox(s4, x + Inches(0.25), Inches(1.85), Inches(3.4), Inches(0.45), phase, 18, True, WHITE, "Georgia")
    add_textbox(s4, x + Inches(0.25), Inches(2.7), Inches(3.4), Inches(0.5), title, 18, True, INK)
    add_textbox(s4, x + Inches(0.25), Inches(3.4), Inches(3.4), Inches(2.0), detail, 13, False, MUTED)

add_textbox(
    s4,
    Inches(0.5),
    Inches(6.5),
    Inches(12),
    Inches(0.4),
    "designtuntas.id  ·  Business Model Canvas  ·  desain hangat untuk pondasi bisnis",
    11,
    False,
    GOLD,
)

out = r"D:\Resume CV & Konsultan Skripsi online\designtuntas-BMC.pptx"
prs.save(out)
print(out)
